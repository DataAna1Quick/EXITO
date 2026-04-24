"""
Convertidor HTML -> PPTX 100% editable.

Reconstruye cada slide del HTML como objetos nativos de PowerPoint:
- cuadros y tarjetas como formas
- imagenes como pictures
- textos como textboxes editables

La fidelidad depende de lo que PowerPoint soporte de CSS, pero el contenido
queda editable directamente dentro del slide.
"""

from __future__ import annotations

import base64
import re
import tempfile
import time
from io import BytesIO
from pathlib import Path
from urllib.parse import unquote, urlparse

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from selenium import webdriver
from selenium.webdriver.chrome.options import Options as ChromeOptions


BASE_DIR = Path(__file__).parent
HTML_FILE = BASE_DIR / "index.html"
OUTPUT_FILE = BASE_DIR / "Operacion_Medellin_100_Editable.pptx"

HTML_SLIDE_WIDTH = 1280
HTML_SLIDE_HEIGHT = 720
PPT_SLIDE_WIDTH_IN = 13.333333
PPT_SLIDE_HEIGHT_IN = 7.5


def px_to_inches(value: float) -> float:
    return float(value) / 96.0


def read_html() -> str:
    return HTML_FILE.read_text(encoding="utf-8")


def parse_html() -> BeautifulSoup:
    return BeautifulSoup(read_html(), "html.parser")


def extract_text(element) -> str:
    if element is None:
        return ""
    return element.get_text(separator=" ", strip=True)


def clean_text(text: str) -> str:
    text = text.replace("\xa0", " ")
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def extract_slides_content(soup: BeautifulSoup) -> list[dict]:
    slides = soup.find_all("div", class_="slide")
    data = []
    for idx, slide in enumerate(slides, start=1):
        data.append(
            {
                "number": idx,
                "id": slide.get("id", f"slide-{idx}"),
                "text": clean_text(extract_text(slide)),
            }
        )
    return data


def build_chrome_driver() -> webdriver.Chrome:
    options = ChromeOptions()
    options.add_argument("--headless=new")
    options.add_argument("--disable-gpu")
    options.add_argument("--hide-scrollbars")
    options.add_argument("--window-size=1600,1000")
    options.add_argument("--force-device-scale-factor=1")
    options.add_argument("--allow-file-access-from-files")
    options.add_argument("--enable-local-file-accesses")
    return webdriver.Chrome(options=options)


def parse_css_color(value: str) -> tuple[RGBColor | None, float]:
    if not value:
        return None, 0.0

    value = value.strip().lower()
    if value in {"transparent", "initial", "inherit"}:
        return None, 0.0

    rgb_match = re.match(
        r"rgba?\((\d+)\s*,\s*(\d+)\s*,\s*(\d+)(?:\s*,\s*([\d.]+))?\)",
        value,
    )
    if rgb_match:
        r, g, b = [int(rgb_match.group(i)) for i in range(1, 4)]
        alpha = float(rgb_match.group(4) or 1.0)
        return RGBColor(r, g, b), alpha

    hex_match = re.match(r"#([0-9a-f]{3}|[0-9a-f]{6})$", value)
    if hex_match:
        hex_value = hex_match.group(1)
        if len(hex_value) == 3:
            hex_value = "".join(ch * 2 for ch in hex_value)
        r = int(hex_value[0:2], 16)
        g = int(hex_value[2:4], 16)
        b = int(hex_value[4:6], 16)
        return RGBColor(r, g, b), 1.0

    return None, 0.0


def normalize_font_name(font_family: str) -> str | None:
    if not font_family:
        return None
    first = font_family.split(",")[0].strip().strip("\"'")
    return first or None


def resolve_image_bytes(src: str) -> bytes | None:
    if not src:
        return None

    if src.startswith("data:"):
        try:
            _, encoded = src.split(",", 1)
            return base64.b64decode(encoded)
        except Exception:
            return None

    parsed = urlparse(src)
    if parsed.scheme == "file":
        path = Path(unquote(parsed.path.lstrip("/")))
        if path.exists():
            return path.read_bytes()
        return None

    local_path = (HTML_FILE.parent / unquote(src)).resolve()
    if local_path.exists():
        return local_path.read_bytes()

    return None


def collect_slide_layout(driver: webdriver.Chrome, slide_index: int, slide_id: str) -> dict:
    driver.execute_script("goTo(arguments[0]);", slide_index)
    time.sleep(0.15)

    return driver.execute_script(
        r"""
        const slideId = arguments[0];

        const rgbaToObj = (value) => {
          if (!value) return null;
          const match = value.match(/rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)/i);
          if (!match) return null;
          return {
            r: Number(match[1]),
            g: Number(match[2]),
            b: Number(match[3]),
            a: match[4] === undefined ? 1 : Number(match[4])
          };
        };

        const isVisible = (el) => {
          const style = window.getComputedStyle(el);
          const rect = el.getBoundingClientRect();
          return style.display !== 'none' &&
            style.visibility !== 'hidden' &&
            Number(style.opacity || 1) > 0.01 &&
            rect.width > 1 && rect.height > 1;
        };

        const slide = document.getElementById(slideId);
        const nav = document.getElementById('nav');
        if (nav) nav.style.display = 'none';

        const slideRect = slide.getBoundingClientRect();
        const slideStyle = window.getComputedStyle(slide);

        const result = {
          backgroundColor: slideStyle.backgroundColor,
          rects: [],
          texts: [],
          images: []
        };

        for (const el of slide.querySelectorAll('*')) {
          if (!isVisible(el)) continue;
          if (el.tagName === 'IMG') continue;

          const style = window.getComputedStyle(el);
          const rect = el.getBoundingClientRect();
          const bg = rgbaToObj(style.backgroundColor);
          const border = rgbaToObj(style.borderTopColor);
          const borderWidth = parseFloat(style.borderTopWidth || '0');

          const hasFill = bg && bg.a > 0.03;
          const hasBorder = border && border.a > 0.03 && borderWidth > 0.2 && style.borderTopStyle !== 'none';
          if (!hasFill && !hasBorder) continue;

          result.rects.push({
            x: rect.left - slideRect.left,
            y: rect.top - slideRect.top,
            w: rect.width,
            h: rect.height,
            bg: style.backgroundColor,
            borderColor: style.borderTopColor,
            borderWidth,
            radius: parseFloat(style.borderTopLeftRadius || '0')
          });
        }

        const walker = document.createTreeWalker(slide, NodeFilter.SHOW_TEXT);
        while (walker.nextNode()) {
          const node = walker.currentNode;
          const raw = node.textContent || '';
          const text = raw.replace(/\s+/g, ' ').trim();
          if (!text) continue;

          const parent = node.parentElement;
          if (!parent || !isVisible(parent)) continue;

          const range = document.createRange();
          range.selectNodeContents(node);
          const rects = Array.from(range.getClientRects()).filter(r => r.width > 0.5 && r.height > 0.5);
          if (!rects.length) continue;

          const left = Math.min(...rects.map(r => r.left));
          const top = Math.min(...rects.map(r => r.top));
          const right = Math.max(...rects.map(r => r.right));
          const bottom = Math.max(...rects.map(r => r.bottom));

          const style = window.getComputedStyle(parent);

          result.texts.push({
            text,
            x: left - slideRect.left,
            y: top - slideRect.top,
            w: right - left,
            h: bottom - top,
            color: style.color,
            fontSize: parseFloat(style.fontSize || '12'),
            fontWeight: String(style.fontWeight || '400'),
            fontStyle: style.fontStyle,
            fontFamily: style.fontFamily,
            textAlign: style.textAlign,
            textTransform: style.textTransform,
            lineHeight: style.lineHeight,
            letterSpacing: style.letterSpacing
          });
        }

        for (const img of slide.querySelectorAll('img')) {
          if (!isVisible(img)) continue;
          const rect = img.getBoundingClientRect();
          result.images.push({
            x: rect.left - slideRect.left,
            y: rect.top - slideRect.top,
            w: rect.width,
            h: rect.height,
            src: img.currentSrc || img.src || ''
          });
        }

        return result;
        """,
        slide_id,
    )


def apply_background(slide, color_value: str) -> None:
    color, alpha = parse_css_color(color_value)
    if color is None:
        return

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color
    if alpha < 1:
        fill.transparency = 1 - alpha


def add_rect_shape(slide, item: dict) -> None:
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if float(item.get("radius", 0) or 0) >= 4
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(px_to_inches(item["x"])),
        Inches(px_to_inches(item["y"])),
        Inches(px_to_inches(item["w"])),
        Inches(px_to_inches(item["h"])),
    )

    bg_color, bg_alpha = parse_css_color(item.get("bg", ""))
    if bg_color is not None and bg_alpha > 0:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if bg_alpha < 1:
            shape.fill.transparency = 1 - bg_alpha
    else:
        shape.fill.background()

    border_color, border_alpha = parse_css_color(item.get("borderColor", ""))
    border_width = float(item.get("borderWidth", 0) or 0)
    if border_color is not None and border_width > 0:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width * 0.75)
        if border_alpha < 1:
            shape.line.transparency = 1 - border_alpha
    else:
        shape.line.fill.background()


def add_textbox(slide, item: dict) -> None:
    textbox = slide.shapes.add_textbox(
        Inches(px_to_inches(item["x"])),
        Inches(px_to_inches(item["y"])),
        Inches(px_to_inches(max(item["w"], 6))),
        Inches(px_to_inches(max(item["h"], item["fontSize"] * 1.4))),
    )
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    paragraph = text_frame.paragraphs[0]
    paragraph.text = transform_text(item["text"], item.get("textTransform"))
    paragraph.alignment = map_alignment(item.get("textAlign"))

    run = paragraph.runs[0]
    run.font.size = Pt(float(item.get("fontSize", 12)) * 0.75)
    run.font.bold = is_bold(item.get("fontWeight", "400"))
    run.font.italic = item.get("fontStyle") == "italic"

    font_name = normalize_font_name(item.get("fontFamily", ""))
    if font_name:
        run.font.name = font_name

    color, alpha = parse_css_color(item.get("color", ""))
    if color is not None:
        run.font.color.rgb = color
        if alpha < 1:
            run.font.fill.transparency = 1 - alpha

    textbox.fill.background()
    textbox.line.fill.background()


def add_image(slide, item: dict) -> None:
    image_bytes = resolve_image_bytes(item.get("src", ""))
    if not image_bytes:
        return

    slide.shapes.add_picture(
        BytesIO(image_bytes),
        Inches(px_to_inches(item["x"])),
        Inches(px_to_inches(item["y"])),
        width=Inches(px_to_inches(item["w"])),
        height=Inches(px_to_inches(item["h"])),
    )


def transform_text(text: str, transform: str | None) -> str:
    if transform == "uppercase":
        return text.upper()
    if transform == "lowercase":
        return text.lower()
    return text


def is_bold(weight: str) -> bool:
    if weight in {"bold", "bolder"}:
        return True
    try:
        return int(weight) >= 600
    except ValueError:
        return False


def map_alignment(value: str | None):
    if value == "center":
        return PP_ALIGN.CENTER
    if value == "right":
        return PP_ALIGN.RIGHT
    if value == "justify":
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


def set_notes_text(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.paragraphs[0].text = text or "Sin contenido extraido."


def create_presentation(slide_models: list[dict]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(PPT_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(PPT_SLIDE_HEIGHT_IN)

    for model in slide_models:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide, model.get("backgroundColor", ""))

        for rect in model.get("rects", []):
            add_rect_shape(slide, rect)

        for image in model.get("images", []):
            add_image(slide, image)

        for text in model.get("texts", []):
            add_textbox(slide, text)

        set_notes_text(slide, model.get("rawText", ""))

    return prs


def main() -> None:
    print("Iniciando conversion HTML -> PPTX 100% editable...")

    soup = parse_html()
    slides = extract_slides_content(soup)
    print(f"Se encontraron {len(slides)} slides")

    driver = build_chrome_driver()
    try:
        driver.get(HTML_FILE.resolve().as_uri())
        time.sleep(1.2)

        if driver.execute_script("return document.fonts && document.fonts.ready ? 1 : 0;"):
            driver.execute_async_script(
                r"""
                const done = arguments[0];
                document.fonts.ready.then(() => done(true)).catch(() => done(false));
                """
            )

        slide_models = []
        for idx, slide_meta in enumerate(slides):
            print(f"Procesando {slide_meta['id']} ({idx + 1}/{len(slides)})...")
            layout = collect_slide_layout(driver, idx, slide_meta["id"])
            layout["id"] = slide_meta["id"]
            layout["rawText"] = slide_meta["text"]
            slide_models.append(layout)

    finally:
        driver.quit()

    print("Construyendo archivo PowerPoint...")
    prs = create_presentation(slide_models)
    prs.save(str(OUTPUT_FILE))
    print(f"Archivo guardado: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
